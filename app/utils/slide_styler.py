from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from typing import Optional

def apply_theme_to_text_frame(text_frame, font: str, color_hex: str):
    rgb = RGBColor.from_string(color_hex.strip("#"))
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(18)
            run.font.color.rgb = rgb

def apply_config_to_presentation(file_path: str, font: str, color_hex: str):
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                apply_theme_to_text_frame(shape.text_frame, font, color_hex)

    prs.save(file_path)
